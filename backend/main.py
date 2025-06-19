from fastapi import FastAPI, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from langchain_openai import ChatOpenAI
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.retrieval_qa.base import RetrievalQA
from langchain_text_splitters import CharacterTextSplitter
import pandas as pd
import os
from dotenv import load_dotenv
import io
import fitz  # PyMuPDF
from bs4 import BeautifulSoup
from pptx import Presentation
import requests

# Load environment variables from .env file
load_dotenv()

app = FastAPI()

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Change to specific origins in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize text splitter
text_splitter = CharacterTextSplitter.from_tiktoken_encoder(
    chunk_size=300,
    chunk_overlap=50
)

# Initialize embedding
embedding = OpenAIEmbeddings(model="text-embedding-3-small")
# embedding = OpenAIEmbeddings(model="gpt-4.1")
VECTOR_DIR = "faiss_index"

# Initialize vector store
try:
    if os.path.exists(VECTOR_DIR):
        vectorstore = FAISS.load_local(
            VECTOR_DIR,
            embeddings=embedding,
            allow_dangerous_deserialization=True
        )
    else:
        vectorstore = FAISS.from_texts(["Initial dummy doc"], embedding)
except Exception as e:
    print(f"Error initializing vector store: {e}")
    vectorstore = FAISS.from_texts(["Initial dummy doc"], embedding)

retriever = vectorstore.as_retriever()
qa_chain = RetrievalQA.from_chain_type(
    llm=ChatOpenAI(model="gpt-3.5-turbo"),
    retriever=retriever
)

class Query(BaseModel):
    text: str

def extract_text_from_url(url: str) -> str:
    """Fetch and extract readable text from a URL"""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, "html.parser")

        # Remove script/style elements
        for tag in soup(["script", "style"]):
            tag.extract()

        text = soup.get_text(separator=" ", strip=True)
        return text
    except Exception as e:
        raise Exception(f"Failed to extract text from URL: {str(e)}")

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(file_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
        return "\n".join(all_text)
    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")

@app.get("/")
async def root():
    return {"message": "FastAPI Chatbot API is running!"}

@app.post("/ask")
async def ask(query: Query):
    try:
        print(f"Question: {query.text}")
        result = qa_chain.invoke({"query": query.text})
        answer = result["result"] if isinstance(result, dict) else str(result)
        print(f"Answer: {answer}")
        return {"response": answer}
    except Exception as e:
        print(f"Failed to process query: {str(e)}")
        return {"error": f"Failed to process query: {str(e)}"}

@app.post("/train")
async def train(query: Query):
    try:
        vectorstore.add_texts([query.text])
        vectorstore.save_local(VECTOR_DIR)
        return {"status": "Training data added successfully."}
    except Exception as e:
        return {"error": f"Failed to add training data: {str(e)}"}

@app.post("/upload_txt")
async def upload_txt_file(file: UploadFile = File(...)):
    if not file.filename.endswith(".txt"):
        return {"error": "Only .txt files are supported."}

    try:
        contents = await file.read()
        text = contents.decode("utf-8")

        if not text.strip():
            return {"error": "The uploaded file is empty."}

        # Split text into chunks
        chunks = text_splitter.split_text(text)
        vectorstore.add_texts(chunks)
        vectorstore.save_local(VECTOR_DIR)

        return {"message": f"File '{file.filename}' uploaded and embedded successfully."}
    except Exception as e:
        return {"error": f"Error processing text file: {str(e)}"}

@app.post("/upload_csv")
async def upload_csv_file(file: UploadFile = File(...)):
    if not file.filename.endswith(".csv"):
        return {"error": "Only CSV files are supported."}

    try:
        contents = await file.read()
        df = pd.read_csv(io.StringIO(contents.decode("utf-8")))

        # Convert each row into a text block
        texts = []
        for _, row in df.iterrows():
            row_text = " | ".join(str(cell) for cell in row if pd.notnull(cell))
            texts.append(row_text)

        if not texts:
            return {"error": "No data found in CSV file."}

        # Add to vector store
        vectorstore.add_texts(texts)
        vectorstore.save_local(VECTOR_DIR)

        return {"message": f"CSV file '{file.filename}' uploaded and trained successfully."}
    except Exception as e:
        return {"error": f"Error processing CSV: {str(e)}"}

@app.post("/upload_pdf")
async def upload_pdf_file(file: UploadFile = File(...)):
    if not file.filename.endswith(".pdf"):
        return {"error": "Only PDF files are supported."}

    try:
        # Read file content
        contents = await file.read()

        # Create temporary file
        temp_path = f"/tmp/{file.filename}"
        with open(temp_path, "wb") as f:
            f.write(contents)

        # Extract text from PDF using PyMuPDF
        doc = fitz.open(temp_path)
        all_text = []
        for page in doc:
            text = page.get_text()
            if text.strip():  # Only add non-empty pages
                all_text.append(text)
        doc.close()

        # Clean up temporary file
        os.remove(temp_path)

        if not all_text:
            return {"error": "No text found in the PDF."}

        # Join all text and split into chunks
        full_text = "\n".join(all_text)
        chunks = text_splitter.split_text(full_text)

        # Add to vector store
        vectorstore.add_texts(chunks)
        vectorstore.save_local(VECTOR_DIR)

        return {"message": f"PDF file '{file.filename}' processed and embedded successfully."}
    except Exception as e:
        return {"error": f"Error processing PDF: {str(e)}"}

@app.post("/train_url")
async def train_from_url(url: str = Form(...)):
    try:
        text = extract_text_from_url(url)
        if not text.strip():
            return {"error": "No text found at the provided URL."}

        # Split text into chunks
        chunks = text_splitter.split_text(text)
        vectorstore.add_texts(chunks)
        vectorstore.save_local(VECTOR_DIR)

        return {"message": f"Content from URL '{url}' embedded and stored successfully."}
    except Exception as e:
        return {"error": f"Failed to process URL: {str(e)}"}

@app.post("/upload_pptx")
async def upload_pptx_file(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        return {"error": "Only .pptx files are supported."}

    try:
        # Read file content
        contents = await file.read()

        # Save the file temporarily
        temp_path = f"/tmp/{file.filename}"
        with open(temp_path, "wb") as f:
            f.write(contents)

        # Extract text
        text = extract_text_from_pptx(temp_path)

        # Clean up temporary file
        os.remove(temp_path)

        if not text.strip():
            return {"error": "No text found in the PowerPoint file."}

        # Split text into chunks
        chunks = text_splitter.split_text(text)

        # Embed and store
        vectorstore.add_texts(chunks)
        vectorstore.save_local(VECTOR_DIR)

        return {"message": f"PowerPoint file '{file.filename}' processed and embedded successfully."}
    except Exception as e:
        return {"error": f"Error processing PowerPoint: {str(e)}"}